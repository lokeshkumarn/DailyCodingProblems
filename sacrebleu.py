
from docx import Document
import nltk
nltk.download('punkt')
from nltk.tokenize import word_tokenizefrom pptx import Presentation
import nltk
nltk.download('punkt')
from nltk.tokenize import word_tokenize
import sacrebleu

# Function to extract all text from PPTX file
def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    full_text = []

    # Extract text from slides
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
        # Extract text from slide notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            for shape in notes_slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)

    # Extract text from slide masters
    for slide_master in prs.slide_masters:
        for shape in slide_master.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
        # Extract text from slide master notes
        if slide_master.has_notes_master:
            notes_master = slide_master.notes_master
            for shape in notes_master.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)

    return '\n'.join(full_text)

# Function to preprocess text
def preprocess_text(text):
    return word_tokenize(text.lower())

# Function to calculate BLEU score
def bleu_score(reference, candidate):
    return sacrebleu.corpus_bleu([candidate], [[reference]]).score

# Function to compare translations
def compare_translations(human_translated_path, machine_translated_path):
    human_text = extract_text_from_pptx(human_translated_path)
    machine_text = extract_text_from_pptx(machine_translated_path)
    
    preprocessed_human_text = preprocess_text(human_text)
    preprocessed_machine_text = preprocess_text(machine_text)
    
    human_text_joined = ' '.join(preprocessed_human_text)
    machine_text_joined = ' '.join(preprocessed_machine_text)
    
    bleu = bleu_score(human_text_joined, machine_text_joined)
    
    return bleu

# Example usage
human_translated_path = 'human_translated.pptx'
machine_translated_path = 'machine_translated.pptx'
bleu_score_value = compare_translations(human_translated_path, machine_translated_path)
print(f"BLEU Score: {bleu_score_value}")

import sacrebleu

# Function to extract all text from DOCX file
def extract_text_from_docx(file_path):
    doc = Document(file_path)
    full_text = []

    # Extract text from paragraphs
    for para in doc.paragraphs:
        full_text.append(para.text)

    # Extract text from tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                full_text.append(cell.text)

    # Extract text from headers and footers
    for section in doc.sections:
        for header in section.header.paragraphs:
            full_text.append(header.text)
        for footer in section.footer.paragraphs:
            full_text.append(footer.text)

    # Extract text from footnotes and endnotes
    if hasattr(doc, 'footnotes'):
        for footnote in doc.footnotes.part.element.getiterator():
            if footnote.tag.endswith('p'):
                full_text.append(footnote.text)
    
    if hasattr(doc, 'endnotes'):
        for endnote in doc.endnotes.part.element.getiterator():
            if endnote.tag.endswith('p'):
                full_text.append(endnote.text)

    return '\n'.join(full_text)

# Function to preprocess text
def preprocess_text(text):
    return word_tokenize(text.lower())

# Function to calculate BLEU score
def bleu_score(reference, candidate):
    return sacrebleu.corpus_bleu([candidate], [[reference]]).score

# Function to compare translations
def compare_translations(human_translated_path, machine_translated_path):
    human_text = extract_text_from_docx(human_translated_path)
    machine_text = extract_text_from_docx(machine_translated_path)
    
    preprocessed_human_text = preprocess_text(human_text)
    preprocessed_machine_text = preprocess_text(machine_text)
    
    human_text_joined = ' '.join(preprocessed_human_text)
    machine_text_joined = ' '.join(preprocessed_machine_text)
    
    bleu = bleu_score(human_text_joined, machine_text_joined)
    
    return bleu

# Example usage
human_translated_path = 'human_translated.docx'
machine_translated_path = 'machine_translated.docx'
bleu_score_value = compare_translations(human_translated_path, machine_translated_path)
print(f"BLEU Score: {bleu_score_value}")

from pptx import Presentation
import nltk
nltk.download('punkt')
from nltk.tokenize import word_tokenize
import sacrebleu

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    full_text = []

    # Extract text from slides
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)

    return '\n'.join(full_text)

# Function to preprocess text
def preprocess_text(text):
    return word_tokenize(text.lower())

# Function to calculate BLEU score
def bleu_score(reference, candidate):
    return sacrebleu.corpus_bleu([candidate], [[reference]]).score

# Function to compare translations
def compare_translations(human_translated_path, machine_translated_path):
    human_text = extract_text_from_pptx(human_translated_path)
    machine_text = extract_text_from_pptx(machine_translated_path)
    
    preprocessed_human_text = preprocess_text(human_text)
    preprocessed_machine_text = preprocess_text(machine_text)
    
    human_text_joined = ' '.join(preprocessed_human_text)
    machine_text_joined = ' '.join(preprocessed_machine_text)
    
    bleu = bleu_score(human_text_joined, machine_text_joined)
    
    return bleu

# Example usage
human_translated_path = 'human_translated.pptx'
machine_translated_path = 'machine_translated.pptx'
bleu_score_value = compare_translations(human_translated_path, machine_translated_path)
print(f"BLEU Score: {bleu_score_value}")

#pip install python-pptx nltk sacrebleu


